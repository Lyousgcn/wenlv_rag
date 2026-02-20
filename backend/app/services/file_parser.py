from pathlib import Path
from typing import Iterable, List, Tuple

import docx
import markdown
from pptx import Presentation
from PyPDF2 import PdfReader


SUPPORTED_EXTENSIONS = {".pdf", ".ppt", ".pptx", ".md", ".markdown", ".doc", ".docx", ".png"}


def extract_text_from_pdf(path: Path) -> str:
    reader = PdfReader(str(path))
    texts: List[str] = []
    for page in reader.pages:
        texts.append(page.extract_text() or "")
    return "\n".join(texts)


def extract_text_from_ppt(path: Path) -> str:
    pres = Presentation(str(path))
    texts: List[str] = []
    for slide in pres.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return "\n".join(texts)


def extract_text_from_md(path: Path) -> str:
    content = path.read_text(encoding="utf-8", errors="ignore")
    html = markdown.markdown(content)
    return html


def extract_text_from_docx(path: Path) -> str:
    document = docx.Document(str(path))
    return "\n".join(p.text for p in document.paragraphs)


def extract_text_from_png(path: Path) -> str:
    return f"图片文件：{path.name}。当前示例环境未集成OCR，仅记录文件名称。"


def extract_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return extract_text_from_pdf(path)
    if suffix in {".ppt", ".pptx"}:
        return extract_text_from_ppt(path)
    if suffix in {".md", ".markdown"}:
        return extract_text_from_md(path)
    if suffix in {".doc", ".docx"}:
        return extract_text_from_docx(path)
    if suffix == ".png":
        return extract_text_from_png(path)
    raise ValueError(f"不支持的文件类型: {suffix}")


def split_text_to_chunks(
    text: str,
    chunk_size: int = 500,
    chunk_overlap: int = 100,
) -> List[str]:
    """简单的按字符切块"""
    if not text:
        return []
    chunks: List[str] = []
    start = 0
    length = len(text)
    while start < length:
        end = min(start + chunk_size, length)
        chunks.append(text[start:end])
        if end == length:
            break
        start = end - chunk_overlap
        if start < 0:
            start = 0
    return chunks


def iter_file_chunks(
    path: Path,
    chunk_size: int,
    chunk_overlap: int,
) -> Iterable[Tuple[int, str]]:
    """生成文件切块序列，返回(索引, 文本)"""
    text = extract_text(path)
    chunks = split_text_to_chunks(text, chunk_size, chunk_overlap)
    for idx, chunk in enumerate(chunks):
        yield idx, chunk

